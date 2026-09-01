import io
import numpy as np
import pandas as pd
import streamlit as st

# python-pptx のインポートチェック
try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# --- ページ設定 ---
st.set_page_config(
    page_title="教育費用シミュレーター", page_icon="🎓", layout="wide"
)

# --- カスタムCSS（Web画面で元スライド7枚目のデザイン・フォントを完全再現） ---
st.markdown(
    """
<style>
    @import url('https://fonts.googleapis.com/css2?family=Noto+Sans+JP:wght@400;500;700;900&family=Poppins:wght@600;700;800&display=swap');
    
    html, body, [class*="css"], .stApp {
        font-family: 'Noto Sans JP', '游ゴシック', 'Yu Gothic', 'メイリオ', sans-serif !important;
    }

    /* KPIカード (元スライド7枚目と同等) */
    .kpi-card {
        background: linear-gradient(135deg, #1e3a8a 0%, #0f172a 100%);
        border-left: 6px solid #f59e0b;
        border-radius: 10px;
        color: #ffffff;
        padding: 22px 30px;
        margin-bottom: 24px;
        box-shadow: 0 4px 16px rgba(15, 23, 42, 0.15);
        display: flex;
        justify-content: space-between;
        align-items: center;
    }
    
    .kpi-title {
        color: #93c5fd;
        font-size: 15px;
        font-weight: 700;
        margin-bottom: 4px;
    }
    
    .kpi-sub {
        color: #cbd5e1;
        font-size: 12px;
    }

    .kpi-value {
        color: #fbbf24;
        font-family: 'Poppins', 'Noto Sans JP', sans-serif;
        font-size: 38px;
        font-weight: 800;
        letter-spacing: -0.5px;
    }

    /* カテゴリタグ (元スライド7枚目のバッジ) */
    .category-tag {
        background-color: #1e3a8a;
        color: #ffffff;
        font-size: 11px;
        font-weight: 700;
        padding: 4px 12px;
        border-radius: 4px;
        display: inline-block;
        margin-bottom: 8px;
        letter-spacing: 0.5px;
    }

    /* コンサル風カスタムテーブル */
    .table-custom-container {
        width: 100%;
        overflow-x: auto;
        margin-top: 10px;
    }

    .table-custom {
        width: 100%;
        border-collapse: separate;
        border-spacing: 0;
        border-radius: 8px;
        overflow: hidden;
        border: 1px solid #cbd5e1;
        box-shadow: 0 2px 8px rgba(0, 0, 0, 0.03);
    }

    .table-custom th {
        background-color: #0f172a;
        color: #f8fafc;
        font-size: 13.5px;
        font-weight: 700;
        padding: 14px 16px;
        text-align: left;
        border-bottom: 2px solid #1e293b;
    }

    .table-custom td {
        background-color: #ffffff;
        border-bottom: 1px solid #e2e8f0;
        color: #334155;
        font-size: 13px;
        padding: 12px 16px;
    }

    .table-custom tr:nth-child(even) td {
        background-color: #f8fafc;
    }

    .table-custom tr:last-child td {
        border-bottom: none;
    }

    .highlight-price {
        color: #1e3a8a;
        font-weight: 700;
    }
    
    .highlight-price-amber {
        color: #d97706;
        font-weight: 800;
    }

    .note-footer {
        font-size: 11.5px;
        color: #64748b;
        margin-top: 10px;
        margin-bottom: 20px;
    }
</style>
""",
    unsafe_allow_html=True,
)

st.title("🎓 受験・教育費用シミュレーター")
st.caption(
    "文部科学省公的データおよび主要塾の実費マスターに基づく学習費用試算"
)


# --- データ読み込み ---
@st.cache_data
def load_data():
    df_master = pd.read_excel(
        "juku_evidence_master_v4.xlsx",
        sheet_name="収集データマスター",
        skiprows=3,
    )
    df_mext = pd.read_excel(
        "mext_juku_expenses_2023.xlsx",
        sheet_name="学年別学習塾費データ",
        skiprows=3,
    )
    df_mext.columns = [str(c).strip() for c in df_mext.columns]
    return df_master, df_mext


try:
    df_master, df_mext = load_data()
except Exception as e:
    st.error(f"データファイルの読み込みに失敗しました。詳細: {e}")
    st.stop()


# --- パワポ生成関数（元スライド7枚目のデザイン・フォントを再現） ---
def create_pptx_download(
    df_res, total_cost, route_label, cost_mode, selected_grades
):
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 画面サイズ
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank_layout)

    # 1. 背景色 (#F1F5F9)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(241, 245, 249)
    bg.line.color.rgb = RGBColor(241, 245, 249)

    # 2. 上部ヘッダーアクセントライン (1E3A8A)
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.12),
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RGBColor(30, 58, 138)
    top_bar.line.color.rgb = RGBColor(30, 58, 138)

    # 3. カテゴリタグ (バッジ: OUTPUT)
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(0.4),
        Inches(2.2),
        Inches(0.35),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = RGBColor(30, 58, 138)
    tag.line.color.rgb = RGBColor(30, 58, 138)
    tf_tag = tag.text_frame
    tf_tag.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = "OUTPUT ： 試算結果モデル"
    p_tag.font.name = "游ゴシック"
    p_tag.font.size = Pt(10)
    p_tag.font.bold = True
    p_tag.font.color.rgb = RGBColor(255, 255, 255)
    p_tag.alignment = PP_ALIGN.CENTER

    # 4. スライドタイトル
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.6)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"個別最適化された通塾費用推移（{route_label}モデル）"
    p.font.name = "游ゴシック"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)

    # 5. KPIカード背景 (1E3A8A -> 0F172A)
    kpi_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(1.5),
        Inches(11.733),
        Inches(1.2),
    )
    kpi_box.fill.solid()
    kpi_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    kpi_box.line.color.rgb = RGBColor(30, 41, 59)

    # KPIカード左アクセント線 (F59E0B)
    kpi_accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8),
        Inches(1.5),
        Inches(0.12),
        Inches(1.2),
    )
    kpi_accent.fill.solid()
    kpi_accent.fill.fore_color.rgb = RGBColor(245, 158, 11)
    kpi_accent.line.color.rgb = RGBColor(245, 158, 11)

    # KPIテキスト
    tx_kpi = slide.shapes.add_textbox(
        Inches(1.1), Inches(1.55), Inches(7.0), Inches(1.1)
    )
    tf_kpi = tx_kpi.text_frame
    p_kpi_1 = tf_kpi.paragraphs[0]
    p_kpi_1.text = (
        f"試算モデル：{route_label} （{len(selected_grades)}年間・{cost_mode}）"
    )
    p_kpi_1.font.name = "游ゴシック"
    p_kpi_1.font.size = Pt(13)
    p_kpi_1.font.bold = True
    p_kpi_1.font.color.rgb = RGBColor(147, 197, 253)

    p_kpi_2 = tf_kpi.add_paragraph()
    p_kpi_2.text = f"※選択された学年: {', '.join(selected_grades)}"
    p_kpi_2.font.name = "游ゴシック"
    p_kpi_2.font.size = Pt(10)
    p_kpi_2.font.color.rgb = RGBColor(203, 213, 225)

    # 金額表示 (右寄せ)
    tx_val = slide.shapes.add_textbox(
        Inches(7.5), Inches(1.6), Inches(4.8), Inches(1.0)
    )
    tf_val = tx_val.text_frame
    p_val = tf_val.paragraphs[0]
    p_val.text = f"¥ {total_cost:,.0f} 円"
    p_val.font.name = "Arial"
    p_val.font.size = Pt(32)
    p_val.font.bold = True
    p_val.font.color.rgb = RGBColor(251, 191, 36)
    p_val.alignment = PP_ALIGN.RIGHT

    # 6. テーブル配置
    rows = len(df_res) + 1
    cols = 5 if "中間ver" in cost_mode else 4

    table_left = Inches(0.8)
    table_top = Inches(2.9)
    table_width = Inches(11.733)
    table_height = Inches(3.8)

    table_shape = slide.shapes.add_table(
        rows, cols, table_left, table_top, table_width, table_height
    )
    table = table_shape.table

    # カラム幅調整
    if cols == 5:
        col_widths = [
            Inches(1.8),
            Inches(2.0),
            Inches(2.0),
            Inches(2.2),
            Inches(3.733),
        ]
    else:
        col_widths = [Inches(1.8), Inches(2.5), Inches(4.5), Inches(2.933)]

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # ヘッダー設定
    headers = (
        [
            "対象学年",
            "文科省平均(円)",
            "大手塾実費(円)",
            "算出費用 [幾何平均] (円)",
            "月額換算目安・時期",
        ]
        if "中間ver" in cost_mode
        else ["対象学年", "算出費用(円)", "月額換算目安・時期", "算出根拠"]
    )

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "游ゴシック"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(248, 250, 252)

    # データ行
    for r_idx, row in df_res.iterrows():
        bg_rgb = (
            RGBColor(248, 250, 252)
            if r_idx % 2 == 1
            else RGBColor(255, 255, 255)
        )

        for c_idx in range(cols):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_rgb
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 学年
        c0 = table.cell(r_idx + 1, 0).text_frame.paragraphs[0]
        c0.text = str(row["学年"])
        c0.font.name = "游ゴシック"
        c0.font.bold = True
        c0.font.size = Pt(11)
        c0.font.color.rgb = RGBColor(15, 23, 42)

        price_color = (
            RGBColor(217, 119, 6)
            if row["is_peak"]
            else RGBColor(30, 58, 138)
        )

        if "中間ver" in cost_mode:
            # 文科省平均
            c1 = table.cell(r_idx + 1, 1).text_frame.paragraphs[0]
            c1.text = f"¥{row['文科省平均']:,.0f}"
            c1.font.name = "游ゴシック"
            c1.font.size = Pt(10.5)
            c1.font.color.rgb = RGBColor(51, 65, 85)

            # 大手実費
            c2 = table.cell(r_idx + 1, 2).text_frame.paragraphs[0]
            c2.text = f"¥{row['大手実費']:,.0f}"
            c2.font.name = "游ゴシック"
            c2.font.size = Pt(10.5)
            c2.font.color.rgb = RGBColor(51, 65, 85)

            # 算出費用
            c3 = table.cell(r_idx + 1, 3).text_frame.paragraphs[0]
            c3.text = f"¥{row['算出費用']:,.0f}"
            c3.font.name = "游ゴシック"
            c3.font.size = Pt(11)
            c3.font.bold = True
            c3.font.color.rgb = price_color

            # メモ
            c4 = table.cell(r_idx + 1, 4).text_frame.paragraphs[0]
            c4.text = str(row["月額目安・メモ"])
            c4.font.name = "游ゴシック"
            c4.font.size = Pt(10)
            c4.font.color.rgb = RGBColor(51, 65, 85)
        else:
            # 算出費用
            c1 = table.cell(r_idx + 1, 1).text_frame.paragraphs[0]
            c1.text = f"¥{row['算出費用']:,.0f}"
            c1.font.name = "游ゴシック"
            c1.font.size = Pt(11)
            c1.font.bold = True
            c1.font.color.rgb = price_color

            # メモ
            c2 = table.cell(r_idx + 1, 2).text_frame.paragraphs[0]
            c2.text = str(row["月額目安・メモ"])
            c2.font.name = "游ゴシック"
            c2.font.size = Pt(10)
            c2.font.color.rgb = RGBColor(51, 65, 85)

            # 根拠
            c3 = table.cell(r_idx + 1, 3).text_frame.paragraphs[0]
            c3.text = str(row["算出根拠"])
            c3.font.name = "游ゴシック"
            c3.font.size = Pt(10)
            c3.font.color.rgb = RGBColor(51, 65, 85)

    # 7. フッター注記
    tx_foot = slide.shapes.add_textbox(
        Inches(0.8), Inches(6.8), Inches(11.733), Inches(0.4)
    )
    tf_foot = tx_foot.text_frame
    p_foot = tf_foot.paragraphs[0]
    p_foot.text = "※文部科学省「令和5年度子供の学習費調査」および主要大手塾マスターデータを幾何平均等により動的試算した結果です。"
    p_foot.font.name = "游ゴシック"
    p_foot.font.size = Pt(9)
    p_foot.font.color.rgb = RGBColor(100, 116, 139)

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io


# --- サイドバー：条件設定 ---
st.sidebar.header("⚙️ 条件設定")
st.sidebar.markdown("**1. 受験機会を選択（複数選択可）**")


def update_grade_selection():
    c_chujuken = st.session_state.get("target_chuju", False)
    c_koujuken = st.session_state.get("target_kouju", False)
    c_daijuken = st.session_state.get("target_daiju", False)

    for g in [
        "小1",
        "小2",
        "小3",
        "小4",
        "小5",
        "小6",
        "中1",
        "中2",
        "中3",
        "高1",
        "高2",
        "高3",
    ]:
        target_state = False
        if c_chujuken and g in ["小4", "小5", "小6"]:
            target_state = True
        if c_koujuken and g in ["中1", "中2", "中3"]:
            target_state = True
        if c_daijuken and g in ["高1", "高2", "高3"]:
            target_state = True

        st.session_state[f"btn_{g}"] = target_state


# 安全なセッション状態の初期化
if "target_chuju" not in st.session_state:
    st.session_state["target_chuju"] = True
if "target_kouju" not in st.session_state:
    st.session_state["target_kouju"] = False
if "target_daiju" not in st.session_state:
    st.session_state["target_daiju"] = False

st.sidebar.checkbox(
    "中学受験", key="target_chuju", on_change=update_grade_selection
)
st.sidebar.checkbox(
    "高校受験", key="target_kouju", on_change=update_grade_selection
)
st.sidebar.checkbox(
    "大学受験", key="target_daiju", on_change=update_grade_selection
)

st.sidebar.divider()

cost_mode = st.sidebar.radio(
    "2. 費用感の算出基準",
    [
        "全平均ver（文科省データ）",
        "中間ver（幾何平均）",
        "大手ver（塾実費データ）",
    ],
    help="中間verは文科省公的平均と大手実費の『幾何平均（Geometric Mean）』を用いて外れ値を平滑化した実質標準値です。",
)

# --- メイン画面：通塾学年の選択 ---
st.subheader("📌 通塾学年の選択（複数選択可能）")
st.write("実際に通塾する学年を選択してください。")

all_grades = [
    "小1",
    "小2",
    "小3",
    "小4",
    "小5",
    "小6",
    "中1",
    "中2",
    "中3",
    "高1",
    "高2",
    "高3",
]

col_elem, col_jhs, col_hs = st.columns(3)
selected_grades = []

with col_elem:
    st.markdown("**【小学生】**")
    for g in all_grades[0:6]:
        if st.checkbox(g, key=f"btn_{g}"):
            selected_grades.append(g)

with col_jhs:
    st.markdown("**【中学生】**")
    for g in all_grades[6:9]:
        if st.checkbox(g, key=f"btn_{g}"):
            selected_grades.append(g)

with col_hs:
    st.markdown("**【高校生】**")
    for g in all_grades[9:12]:
        if st.checkbox(g, key=f"btn_{g}"):
            selected_grades.append(g)

st.divider()

# --- 計算・動的レンダリング部 ---
if not selected_grades:
    st.warning(
        "学年が1つも選択されていません。上部のチェックボックスから通塾学年を選択してください。"
    )
else:
    results = []
    mext_spend_col = df_mext.columns[4]

    for grade in selected_grades:
        mext_row = df_mext[df_mext["学年"] == grade]
        mext_val = (
            mext_row[mext_spend_col].values[0] if not mext_row.empty else 0
        )

        if grade in ["小1", "小2", "小3"]:
            route_filter = "中学受験（早期）"
        elif grade in ["小4", "小5", "小6"]:
            # 中学受験チェックが入っていない場合は「基礎固め・高校受験準備」を適用
            route_filter = (
                "中学受験"
                if st.session_state.get("target_chuju", False)
                else "基礎固め・高校受験準備"
            )
        elif grade in ["中1", "中2", "中3"]:
            route_filter = "高校受験"
        elif grade in ["高1", "高2", "高3"]:
            route_filter = "大学受験"

        master_filtered = df_master[
            (df_master["対象ルート"] == route_filter)
            & (df_master["対象学年"] == grade)
        ]

        if not master_filtered.empty:
            juku_val = master_filtered["年間総計(円)"].mean()
            juku_name = ", ".join(master_filtered["塾・情報元名称"].unique())
        else:
            juku_val = mext_val
            juku_name = "参考推計値"

        if "全平均" in cost_mode:
            final_cost = mext_val
            calc_note = "文科省：公立通塾者平均"
        elif "大手" in cost_mode:
            final_cost = juku_val
            calc_note = f"大手実費 ({juku_name})"
        else:  # 中間ver（幾何平均）
            if mext_val > 0 and juku_val > 0:
                final_cost = np.sqrt(mext_val * juku_val)
            else:
                final_cost = max(mext_val, juku_val)
            calc_note = "幾何平均 (文科省 × 大手)"

        # ピーク期判定
        monthly_cost = int(final_cost / 12)
        if grade in ["小6", "中3", "高3"]:
            peak_note = (
                f"約 ¥{monthly_cost:,.0f}/月（受験本番：志望校特訓・直前対策等）"
            )
            is_peak = True
        elif grade in ["小5", "中2", "高2"]:
            peak_note = f"約 ¥{monthly_cost:,.0f}/月（学習本格化・夏冬講習拡大）"
            is_peak = False
        else:
            peak_note = f"約 ¥{monthly_cost:,.0f}/月（基礎固め・初期費用・季節講習）"
            is_peak = False

        results.append(
            {
                "学年": grade,
                "文科省平均": int(mext_val),
                "大手実費": int(juku_val),
                "算出費用": int(final_cost),
                "月額目安・メモ": peak_note,
                "算出根拠": calc_note,
                "is_peak": is_peak,
            }
        )

    df_res = pd.DataFrame(results)
    total_cost = df_res["算出費用"].sum()

    st.subheader("📊 シミュレーション結果（個別最適化モデル）")

    routes_selected = []
    if st.session_state.get("target_chuju"):
        routes_selected.append("中学受験")
    if st.session_state.get("target_kouju"):
        routes_selected.append("高校受験")
    if st.session_state.get("target_daiju"):
        routes_selected.append("大学受験")
    route_label = (
        "・".join(routes_selected) if routes_selected else "カスタム選択"
    )

    # 1. KPIカード
    st.markdown(
        f"""<div class="kpi-card"><div><div class="kpi-title">試算モデル：{route_label} （{len(selected_grades)}年間・{cost_mode}）</div><div class="kpi-sub">※選択された学年：{', '.join(selected_grades)}</div></div><div class="kpi-value">¥ {total_cost:,.0f} 円</div></div>""",
        unsafe_allow_html=True,
    )

    # 2. テーブル
    st.markdown("#### 📑 学年別費用推移内訳")

    if "中間ver" in cost_mode:
        table_html = """<div class="table-custom-container"><table class="table-custom"><thead><tr><th>対象学年</th><th>文科省平均(円)</th><th>大手塾実費(円)</th><th>算出費用 [幾何平均] (円)</th><th>月額換算目安・費用ピーク時期</th></tr></thead><tbody>"""
        for _, row in df_res.iterrows():
            price_class = (
                "highlight-price-amber"
                if row["is_peak"]
                else "highlight-price"
            )
            table_html += f"""<tr><td><strong>{row['学年']}</strong></td><td>¥{row['文科省平均']:,.0f}</td><td>¥{row['大手実費']:,.0f}</td><td><span class="{price_class}">¥{row['算出費用']:,.0f}</span></td><td>{row['月額目安・メモ']}</td></tr>"""
    else:
        val_title = (
            "文科省平均(円)" if "全平均" in cost_mode else "大手塾実費(円)"
        )
        table_html = f"""<div class="table-custom-container"><table class="table-custom"><thead><tr><th>対象学年</th><th>{val_title}</th><th>月額換算目安・費用ピーク時期</th><th>算出根拠</th></tr></thead><tbody>"""
        for _, row in df_res.iterrows():
            price_class = (
                "highlight-price-amber"
                if row["is_peak"]
                else "highlight-price"
            )
            table_html += f"""<tr><td><strong>{row['学年']}</strong></td><td><span class="{price_class}">¥{row['算出費用']:,.0f}</span></td><td>{row['月額目安・メモ']}</td><td>{row['算出根拠']}</td></tr>"""

    table_html += """</tbody></table></div><div class="note-footer">※上記数値は文部科学省「令和5年度子供の学習費調査」および主要塾マスターデータを基に自動動的試算されています。</div>"""

    st.markdown(table_html, unsafe_allow_html=True)

    # 3. PPTX ダウンロードボタン
    if HAS_PPTX:
        pptx_file = create_pptx_download(
            df_res, total_cost, route_label, cost_mode, selected_grades
        )
        st.download_button(
            label="📥 このシミュレーション結果をパワーポイント（.pptx）でダウンロード",
            data=pptx_file,
            file_name="教育費用シミュレーション結果.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )
    else:
        st.info(
            "💡 `pip install python-pptx` を実行すると、パワポファイルのダウンロード機能が有効化されます。"
        )
